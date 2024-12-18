import os
import sys
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle


def extract_text_from_pdf(pdf_file):
    text = ""
    with fitz.open(pdf_file) as pdf:
        for page in pdf:
            text += page.get_text() + "\n"
    return text.strip()


def extract_text_from_docx(docx_file):
    doc = Document(docx_file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text).strip()


def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()


def combine_texts_to_pdf(folder_path, output_pdf_path):
    combined_text = ""

    # Process PDF files
    pdf_files = [f for f in os.listdir(folder_path) if f.lower().endswith('.pdf')]
    for pdf_file in pdf_files:
        full_path = os.path.join(folder_path, pdf_file)
        print(f"Extracting text from {full_path}")
        combined_text += extract_text_from_pdf(full_path) + "\n\n"

    # Process DOCX files
    docx_files = [f for f in os.listdir(folder_path) if f.lower().endswith('.docx')]
    for docx_file in docx_files:
        full_path = os.path.join(folder_path, docx_file)
        print(f"Extracting text from {full_path}")
        combined_text += extract_text_from_docx(full_path) + "\n\n"

    # Process PPTX files
    pptx_files = [f for f in os.listdir(folder_path) if f.lower().endswith('.pptx')]
    for pptx_file in pptx_files:
        full_path = os.path.join(folder_path, pptx_file)
        print(f"Extracting text from {full_path}")
        combined_text += extract_text_from_pptx(full_path) + "\n\n"

    # Create a PDF with the combined text using reportlab
    doc = SimpleDocTemplate(output_pdf_path, pagesize=letter,
                            leftMargin=40, rightMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()

    # Define a custom style with double spacing
    normal_style = ParagraphStyle(name='Normal', parent=styles['Normal'], fontSize=12, leading=16)

    story = []

    # Split the combined text into paragraphs
    for paragraph in combined_text.split('\n\n'):
        p = Paragraph(paragraph, normal_style)
        story.append(p)
        story.append(Spacer(1, 12))  # Add a small space between paragraphs

    doc.build(story)


def main():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.join(script_dir, '..')  # Move up one directory to the project root
    lectures_folder = os.path.join(project_root, "Lectures")

    if not os.path.isdir(lectures_folder):
        print(f"The 'Lectures' folder does not exist in {project_root}. Please create it.")
        sys.exit(1)

    for lecture_subfolder in os.listdir(lectures_folder):
        full_lecture_path = os.path.join(lectures_folder, lecture_subfolder)

        if os.path.isdir(full_lecture_path):
            print(f"Processing lecture folder: {full_lecture_path}")
            output_pdf_path = os.path.join(project_root, f"{lecture_subfolder}.pdf")
            combine_texts_to_pdf(full_lecture_path, output_pdf_path)
            print(f"Combined text saved to {output_pdf_path}")


if __name__ == "__main__":
    main()
